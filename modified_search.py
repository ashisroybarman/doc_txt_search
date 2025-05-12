import os
import threading
import configparser
import zipfile
import textract
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext
from tkinter import ttk  # for Progressbar
import tkinter.font as tkfont
from difflib import SequenceMatcher

# -----------------------------------------------------------------------------
# Helper: Normalize file paths
# -----------------------------------------------------------------------------
def normalize_filepath(filepath):
    """
    Returns an absolute, normalized file path as a raw string.
    This ensures that the path uses a consistent Windows-style separator.
    """
    return r"{}".format(os.path.abspath(filepath))

# -----------------------------------------------------------------------------
# Helper functions for extracting text from supported file types.
# -----------------------------------------------------------------------------
def extract_text_from_doc(filepath):
    """
    Extract text from a .doc file using COM automation.
    """
    try:
        import win32com.client
    except ImportError:
        raise Exception("pywin32 is required for processing .doc files. Please install it.")
    try:
        full_path = normalize_filepath(filepath)
        word = win32com.client.Dispatch("Word.Application")
        # Open the document in read-only mode.
        doc = word.Documents.Open(full_path, ReadOnly=True)
        text = doc.Content.Text
        doc.Close(False)
        word.Quit()
        return text
    except Exception as e:
        raise Exception(f"Error extracting DOC using COM: {e}")

def extract_text_from_docx(filepath):
    """
    Extract text from a .docx file using python-docx.
    """
    try:
        from docx import Document
    except ImportError:
        raise Exception("Please install python-docx to process .docx files.")
    try:
        full_path = normalize_filepath(filepath)
        doc = Document(full_path)
        fullText = [para.text for para in doc.paragraphs]
        return "\n".join(fullText)
    except Exception as e:
        raise Exception(f"Error extracting DOCX: {e}")

def extract_text_from_xlsx(filepath):
    """
    Extract text from a .xlsx file using openpyxl.
    First, check that the file is a valid zip archive.
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise Exception("Please install openpyxl to process .xlsx files.")
    try:
        full_path = normalize_filepath(filepath)
        if not zipfile.is_zipfile(full_path):
            raise Exception("File is not a valid .xlsx zip archive.")
        wb = load_workbook(full_path, data_only=True)
        content = ""
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell is not None)
                content += row_text.strip() + "\n"
        return content
    except Exception as e:
        try:
            full_path = normalize_filepath(filepath)
            content_bytes = textract.process(full_path)
            return content_bytes.decode('utf-8', errors="ignore")
        except Exception as e2:
            raise Exception(f"Error extracting XLSX: {e} / Fallback: {e2}")

def extract_text_from_xls(filepath):
    """
    Extract text from a .xls file using xlrd.
    """
    try:
        import xlrd
    except ImportError:
        raise Exception("Please install xlrd to process .xls files.")
    try:
        full_path = normalize_filepath(filepath)
        book = xlrd.open_workbook(full_path)
        content = ""
        for sheet in book.sheets():
            for i in range(sheet.nrows):
                row = sheet.row_values(i)
                row_text = " ".join(str(cell) for cell in row if cell)
                content += row_text.strip() + "\n"
        return content
    except Exception as e:
        raise Exception(f"Error extracting XLS: {e}")

def extract_text_from_ppt(filepath):
    """
    Extract text from a .ppt file using PowerPoint COM automation.
    """
    try:
        import win32com.client
    except ImportError:
        raise Exception("pywin32 is required for processing .ppt files. Please install it.")
    try:
        full_path = normalize_filepath(filepath)
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        # Open presentation without a visible window.
        presentation = ppt.Presentations.Open(full_path, WithWindow=False)
        text = ""
        for slide in presentation.Slides:
            for shape in slide.Shapes:
                if shape.HasTextFrame and shape.TextFrame.HasText:
                    text += shape.TextFrame.TextRange.Text + "\n"
        presentation.Close()
        ppt.Quit()
        return text
    except Exception as e:
        raise Exception(f"Error extracting PPT using COM: {e}")

def extract_text_from_pptx(filepath):
    """
    Extract text from a .pptx file using python-pptx.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise Exception("Please install python-pptx to process .pptx files.")
    try:
        full_path = normalize_filepath(filepath)
        prs = Presentation(full_path)
        content = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content += shape.text + "\n"
        return content
    except Exception as e:
        raise Exception(f"Error extracting PPTX: {e}")

def extract_text_from_pdf(filepath):
    """
    Extract text from a .pdf file using PyPDF2.
    """
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        raise Exception("Please install PyPDF2 to process .pdf files.")
    try:
        full_path = normalize_filepath(filepath)
        reader = PdfReader(full_path)
        content = ""
        for page in reader.pages:
            content += page.extract_text() + "\n"
        return content
    except Exception as e:
        raise Exception(f"Error extracting PDF: {e}")

def extract_text_from_txt(filepath):
    """
    Extract text from a .txt file.
    """
    try:
        with open(filepath, "r", encoding="utf-8") as file:
            return file.read()
    except Exception as e:
        raise Exception(f"Error extracting TXT: {e}")

# -----------------------------------------------------------------------------
# (Fuzzy) Matching function.
# -----------------------------------------------------------------------------
def combined_match_document(search_str, doc_text, gap_tol=5, context_chars=20, threshold=50, min_match_chars=3):
    """
    Searches for exact occurrences (100% match) of search_str in doc_text.
    If no exact match is found, uses SequenceMatcher for a fuzzy match.
    Returns a list of tuples: [(snippet, percentage), ...]
    """
    trimmed_query = search_str.strip()
    exact_matches = []
    start = 0
    while True:
        idx = doc_text.find(trimmed_query, start)
        if idx == -1:
            break
        exact_matches.append(idx)
        start = idx + len(trimmed_query)
    
    if exact_matches:
        results = []
        for idx in exact_matches:
            snippet_start = max(idx - context_chars, 0)
            snippet_end = min(idx + len(trimmed_query) + context_chars, len(doc_text))
            snippet = doc_text[snippet_start:snippet_end].strip()
            print(f"Exact match for '{search_str}' at index {idx}.")
            results.append((snippet, 100))
        return results

    matcher = SequenceMatcher(None, search_str, doc_text)
    blocks = matcher.get_matching_blocks()
    blocks = [b for b in blocks if b[2] > 0]
    if not blocks:
        return None
    groups = []
    current_group = [blocks[0]]
    for b in blocks[1:]:
        if (b[1] - (current_group[-1][1] + current_group[-1][2])) <= gap_tol:
            current_group.append(b)
        else:
            groups.append(current_group)
            current_group = [b]
    groups.append(current_group)
    results = []
    for group in groups:
        group_start = group[0][1]
        group_end = group[-1][1] + group[-1][2]
        merged_block = doc_text[group_start:group_end]
        trimmed_block = merged_block.strip()
        if len(trimmed_block) < min_match_chars:
            continue
        ratio = SequenceMatcher(None, trimmed_query, trimmed_block).ratio() * 100
        if trimmed_query == trimmed_block:
            ratio = 100
        if ratio >= threshold:
            print(f"Final merged block for '{search_str}':")
            print(repr(merged_block))
            snippet_start = max(group_start - context_chars, 0)
            snippet_end = min(group_end + context_chars, len(doc_text))
            snippet = doc_text[snippet_start:snippet_end].strip()
            results.append((snippet, int(ratio)))
    if results:
        results.sort(key=lambda x: x[1], reverse=True)
        return results
    else:
        return None

# -----------------------------------------------------------------------------
# Tkinter GUI Application with INI File Support and Password-Protected File Detection
# -----------------------------------------------------------------------------
class SearchApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("Document Search App")
        self.window_width, self.window_height = self.winfo_screenwidth() // 3 * 2,  self.winfo_screenheight() // 10 * 8
        self.geometry(f"{self.window_width}x{self.window_height}+10+20")
        # self.geometry("900x700")
        
        # Main frame on the right.
        self.main_frame = tk.Frame(self)
        self.main_frame.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        # Variables for file paths.
        self.search_strings_file = tk.StringVar()
        self.search_target_path = tk.StringVar()
        
        # Allowed file types.
        self.filetypes = {
            ".doc": tk.BooleanVar(value=True),
            ".docx": tk.BooleanVar(value=True),
            ".xls": tk.BooleanVar(value=True),
            ".xlsx": tk.BooleanVar(value=True),
            ".ppt": tk.BooleanVar(value=True),
            ".pptx": tk.BooleanVar(value=True),
            ".pdf": tk.BooleanVar(value=True),  # Added PDF
            ".txt": tk.BooleanVar(value=True)  # Added TXT
        }
        
        # Initialize file type checkboxes to be deselected.
        for ext in self.filetypes:
            self.filetypes[ext].set(False)
        
        # Mapping for clickable file names.
        self.file_tag_map = {}
        self.file_tag_counter = 0
        
        # Bold font.
        self.bold_font = tkfont.Font(family="Helvetica", size=10, weight="bold")
        
        # Update the config file path to "C:/search_app/config.ini".
        self.config_file_path = "C:/search_app/config.ini"
        
        # Load saved configuration.
        self.load_config()
        self.protocol("WM_DELETE_WINDOW", self.on_closing)
        
        self.create_widgets()
    
    def create_widgets(self):
        # Controls frame.
        self.controls_frame = tk.Frame(self.main_frame)
        self.controls_frame.pack(fill=tk.X, anchor="w")
        
        # Search Strings File and "Perform Search" on one line.
        frame_top = tk.Frame(self.controls_frame, width=self.window_width // 10 * 5)
        frame_top.pack(fill=tk.X, anchor="w")
        tk.Label(frame_top, text="Query File:", width = 15).pack(side=tk.LEFT)
        tk.Entry(frame_top, textvariable=self.search_strings_file, width=self.window_width // 10).pack(side=tk.LEFT, padx=5)
        tk.Button(frame_top, text="Browse File", command=self.browse_search_strings).pack(side=tk.LEFT)
        tk.Button(frame_top, text="Perform Search", command=self.start_search, background='lightgreen').pack(side=tk.RIGHT, padx=5)
        
        # Search Target selection.
        frame_target = tk.Frame(self.controls_frame)
        frame_target.pack(fill=tk.X, anchor="w")
        tk.Label(frame_target, text="Target (File/Dir):", width=15).pack(side=tk.LEFT)
        tk.Entry(frame_target, textvariable=self.search_target_path, width=self.window_width // 10).pack(side=tk.LEFT, padx=5)
        tk.Button(frame_target, text="Browse File", command=self.browse_target_file).pack(side=tk.LEFT)
        tk.Button(frame_target, text="Browse Directory", command=self.browse_target_directory).pack(side=tk.LEFT)

        # Add a stop button.
        self.stop_search = False  # Initialize a flag to control the search process.
        tk.Button(frame_target, text="Stop Search", command=self.stop_process, background='red').pack()
        
        # File Types selection (sorted alphabetically).
        frame_types = tk.LabelFrame(self.controls_frame, text="File Types to Search")
        frame_types.pack(fill=tk.X, anchor="e", pady=5)
        for ext, var in sorted(self.filetypes.items()):
            tk.Checkbutton(frame_types, text=ext, variable=var).pack(side=tk.LEFT, padx=5)
        
        # Add a label to display processed files count.
        self.processed_files_label = tk.Label(self.controls_frame, text="Information will be displayed here.", anchor="w")
        self.processed_files_label.pack()
        
        # Results text widget.
        self.txt_results = scrolledtext.ScrolledText(self.main_frame, wrap=tk.WORD, height=25)
        self.txt_results.pack(fill=tk.BOTH, expand=True, pady=5)
    
    def browse_search_strings(self):
        path = filedialog.askopenfilename(title="Select Search Strings File",
                                          filetypes=[("Text Files", "*.txt"), ("All Files", "*.*")])
        if path:
            self.search_strings_file.set(path)
    
    def browse_target_file(self):
        path = filedialog.askopenfilename(title="Select Target File",
                                          filetypes=[("Supported Files", "*.doc;*.docx;*.xls;*.xlsx;*.ppt;*.pptx;*.pdf;*.txt")])
        if path:
            self.search_target_path.set(path)
    
    def browse_target_directory(self):
        path = filedialog.askdirectory(title="Select Target Directory")
        if path:
            self.search_target_path.set(path)
    
    def start_search(self):
        thread = threading.Thread(target=self.perform_search)
        thread.start()
    
    def stop_process(self):
        """
        Set the stop_search flag to True to stop the search process.
        """
        self.stop_search = True

    def perform_search(self):
        self.stop_search = False  # Reset the stop flag to allow the search to restart.
        self.txt_results.delete("1.0", tk.END)
        self.file_tag_map.clear()
        self.file_tag_counter = 0

        search_file = self.search_strings_file.get()
        target_path = self.search_target_path.get()
        if not search_file or not os.path.exists(search_file):
            messagebox.showerror("Error", "Search strings file not selected or does not exist.")
            return
        if not target_path or not os.path.exists(target_path):
            messagebox.showerror("Error", "Search target not selected or does not exist.")
            return

        files_to_search = []
        allowed_exts = self.get_selected_file_types()
        if os.path.isdir(target_path):
            base_dir = target_path
            for root, dirs, files in os.walk(target_path):
                for file in files:
                    if file.startswith("~$"):
                        continue
                    ext = os.path.splitext(file)[1].lower()
                    if ext in allowed_exts:
                        files_to_search.append(os.path.join(root, file))
        else:
            ext = os.path.splitext(target_path)[1].lower()
            if ext in allowed_exts and not os.path.basename(target_path).startswith("~$"):
                files_to_search.append(target_path)
            base_dir = os.path.dirname(target_path)

        if not files_to_search:
            self.txt_results.insert(tk.END, "No files found with the selected file types.\n")
            return

        total_files = len(files_to_search)
        self.processed_files_label.config(text=f"0 out of {total_files} files processed")

        try:
            with open(search_file, "r", encoding="utf-8") as f:
                search_strings = [line.strip() for line in f if line.strip()]
        except Exception as e:
            messagebox.showerror("Error", f"Failed to read search strings file: {e}")
            return

        grouped_results = {s: {} for s in search_strings}
        error_files = []  # To hold files that are password protected.

        for idx, file in enumerate(files_to_search):
            if self.stop_search:  # Check if the stop button was pressed.
                self.txt_results.insert(tk.END, "\nSearch stopped by user.\n")
                self.processed_files_label.config(text="Search stopped.")
                return

            ext = os.path.splitext(file)[1].lower()
            try:
                if ext == ".doc":
                    content = extract_text_from_doc(file)
                elif ext == ".docx":
                    content = extract_text_from_docx(file)
                elif ext == ".ppt":
                    content = extract_text_from_ppt(file)
                elif ext == ".pptx":
                    content = extract_text_from_pptx(file)
                elif ext == ".xlsx":
                    content = extract_text_from_xlsx(file)
                elif ext == ".xls":
                    content = extract_text_from_xls(file)
                elif ext == ".pdf":
                    content = extract_text_from_pdf(file)  # Added PDF handling
                elif ext == ".txt":
                    content = extract_text_from_txt(file)  # Added TXT handling
                else:
                    continue
            except Exception as e:
                err_msg = str(e).lower()
                if "password" in err_msg or "protected" in err_msg:
                    relative_path = os.path.relpath(file, base_dir)
                    error_files.append(relative_path)
                    self.processed_files_label.config(text=f"{idx + 1} out of {total_files} files processed ({relative_path})")
                    continue
                else:
                    self.txt_results.insert(tk.END, f"Error processing file {file}:\n  {e}\n")
                    self.processed_files_label.config(text=f"{idx + 1} out of {total_files} files processed ({relative_path})")
                    continue

            relative_path = os.path.relpath(file, base_dir)
            for s in search_strings:
                matches = combined_match_document(s, content, gap_tol=5, context_chars=20, threshold=50)
                if matches:
                    grouped_results[s].setdefault(relative_path, []).extend(matches)
            self.processed_files_label.config(text=f"{idx + 1} out of {total_files} files processed ({relative_path})")

        for s in search_strings:
            self.txt_results.insert(tk.END, f"'{s}'\n")
            group = grouped_results[s]
            if group:
                for file_rel, match_list in group.items():
                    tag_name = f"fname_{self.file_tag_counter}"
                    self.file_tag_map[tag_name] = os.path.join(base_dir, file_rel)
                    self.file_tag_counter += 1
                    self.txt_results.insert(tk.END, f"\t-{file_rel}\n", tag_name)
                    self.txt_results.tag_configure(tag_name, font=self.bold_font, foreground="blue")
                    self.txt_results.tag_bind(tag_name, "<Double-Button-1>", self.on_file_double_click)
                    for snippet, perc in match_list:
                        self.txt_results.insert(tk.END, f"\t\tMatch = {perc}%, Snippet: {snippet}\n")
            else:
                self.txt_results.insert(tk.END, "\t-NO MATCH FOUND (Match < 50%)\n")
            self.txt_results.insert(tk.END, "-" * 60 + "\n\n")

        if error_files:
            self.txt_results.insert(tk.END, "\nThe following files were skipped because they are password protected:\n")
            for err_file in error_files:
                self.txt_results.insert(tk.END, f"\t- The file {err_file} is password protected.\n")
        self.txt_results.insert(tk.END, "\nSearch completed.\n")
        self.processed_files_label.config(text="Search completed.")
    
    def on_file_double_click(self, event):
        index = self.txt_results.index(f"@{event.x},{event.y}")
        tags = self.txt_results.tag_names(index)
        for tag in tags:
            if tag.startswith("fname_") and tag in self.file_tag_map:
                filepath = self.file_tag_map[tag]
                try:
                    os.startfile(filepath)
                except Exception as e:
                    messagebox.showerror("Error", f"Could not open file:\n{filepath}\n{e}")
                break
    
    def get_selected_file_types(self):
        return sorted([ext for ext, var in self.filetypes.items() if var.get()])
    
    # Update the load_config method to use the new path.
    def load_config(self):
        self.config = configparser.ConfigParser()
        if os.path.exists(self.config_file_path):
            self.config.read(self.config_file_path)
            if "Paths" in self.config:
                if "search_strings_file" in self.config["Paths"]:
                    self.search_strings_file.set(self.config["Paths"]["search_strings_file"])
                if "search_target_path" in self.config["Paths"]:
                    self.search_target_path.set(self.config["Paths"]["search_target_path"])
    
    # Update the save_config method to use the new path.
    def save_config(self):
        self.config["Paths"] = {
            "search_strings_file": self.search_strings_file.get(),
            "search_target_path": self.search_target_path.get()
        }
        os.makedirs(os.path.dirname(self.config_file_path), exist_ok=True)
        with open(self.config_file_path, "w") as configfile:
            self.config.write(configfile)
    
    def on_closing(self):
        self.save_config()
        self.destroy()

# -----------------------------------------------------------------------------
# Run the Application.
# -----------------------------------------------------------------------------
if __name__ == "__main__":
    app = SearchApp()
    app.mainloop()
